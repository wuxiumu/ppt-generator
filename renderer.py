"""PPTX rendering engine — 8 layout types, 5 color palettes."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 16:9 slide dimensions ──────────────────────────────────
SW = 13.333  # slide width  in inches
SH = 7.5     # slide height in inches


# ── Color Palettes ─────────────────────────────────────────
PALETTES = {
    "tech": {
        "bg": "#FAFBFC", "surface": "#EEF0F5",
        "text": "#1A1A2E", "text2": "#4A4A6A", "text3": "#8888AA",
        "accent1": "#6C63FF", "accent2": "#00D2FF", "accent3": "#FF6B6B", "accent4": "#2ED573",
        "on_accent": "#FFFFFF", "divider": "#E2E8F0", "code_bg": "#1E1E2E",
    },
    "warm": {
        "bg": "#FFF8F0", "surface": "#FFF0E0",
        "text": "#2D1B00", "text2": "#6B4E37", "text3": "#A08B7A",
        "accent1": "#C67B5C", "accent2": "#E87461", "accent3": "#4A7C59", "accent4": "#F4C542",
        "on_accent": "#FFFFFF", "divider": "#E8D5C4", "code_bg": "#2D1B00",
    },
    "corporate": {
        "bg": "#FFFFFF", "surface": "#F7FAFC",
        "text": "#1A202C", "text2": "#4A5568", "text3": "#A0AEC0",
        "accent1": "#2B6CB0", "accent2": "#ED8936", "accent3": "#38A169", "accent4": "#E53E3E",
        "on_accent": "#FFFFFF", "divider": "#E2E8F0", "code_bg": "#1A202C",
    },
    "personal": {
        "bg": "#FFFDF7", "surface": "#FAF5FF",
        "text": "#2D3748", "text2": "#718096", "text3": "#A0AEC0",
        "accent1": "#9F7AEA", "accent2": "#F687B3", "accent3": "#68D391", "accent4": "#FBD38D",
        "on_accent": "#FFFFFF", "divider": "#E9D8FD", "code_bg": "#2D3748",
    },
    "dark": {
        "bg": "#0F172A", "surface": "#1E293B",
        "text": "#F8FAFC", "text2": "#CBD5E1", "text3": "#64748B",
        "accent1": "#818CF8", "accent2": "#34D399", "accent3": "#FB923C", "accent4": "#F472B6",
        "on_accent": "#0F172A", "divider": "#334155", "code_bg": "#020617",
    },
}


# ── Layout Specifications (fraction of slide W×H) ────────
# pos = (x, y), size = (w, h), all 0.0–1.0
LAYOUTS = {
    "title": {
        "title":     {"p": (0.08, 0.30), "s": (0.84, 0.16)},
        "subtitle":  {"p": (0.12, 0.50), "s": (0.76, 0.10)},
    },
    "section": {
        "number":    {"p": (0.10, 0.22), "s": (0.18, 0.28)},
        "title":     {"p": (0.10, 0.50), "s": (0.80, 0.14)},
        "subtitle":  {"p": (0.10, 0.66), "s": (0.80, 0.08)},
    },
    "bullets": {
        "title":  {"p": (0.08, 0.06), "s": (0.82, 0.12)},
        "body":   {"p": (0.12, 0.22), "s": (0.78, 0.66)},
    },
    "statement": {
        "quote":  {"p": (0.10, 0.22), "s": (0.80, 0.38)},
        "attr":   {"p": (0.10, 0.66), "s": (0.80, 0.06)},
    },
    "comparison": {
        "title":  {"p": (0.08, 0.06), "s": (0.84, 0.10)},
        "left":   {"p": (0.08, 0.20), "s": (0.40, 0.68)},
        "right":  {"p": (0.52, 0.20), "s": (0.40, 0.68)},
    },
    "big_number": {
        "number": {"p": (0.08, 0.12), "s": (0.84, 0.40)},
        "label":  {"p": (0.08, 0.55), "s": (0.84, 0.10)},
        "context":{"p": (0.14, 0.68), "s": (0.72, 0.08)},
    },
    "code": {
        "title":  {"p": (0.08, 0.06), "s": (0.84, 0.10)},
        "code":   {"p": (0.08, 0.20), "s": (0.50, 0.66)},
        "notes":  {"p": (0.62, 0.20), "s": (0.32, 0.66)},
    },
    "end": {
        "thanks": {"p": (0.10, 0.30), "s": (0.80, 0.16)},
        "contact":{"p": (0.20, 0.52), "s": (0.60, 0.14)},
    },
}


class Renderer:
    """Turns structured slide data into a .pptx file."""

    def render(self, slides: list[dict], palette_name: str,
               topic: str = "", out: str = "output.pptx") -> str:
        self.prs = Presentation()
        self.prs.slide_width  = Inches(SW)
        self.prs.slide_height = Inches(SH)
        self.pal = PALETTES.get(palette_name, PALETTES["tech"])
        n = len(slides)

        for i, s in enumerate(slides):
            layout = s.get("layout", "bullets")
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank

            # background
            self._bg(slide, self.pal["bg"])

            # dispatch
            fn = {
                "title":      self._do_title,
                "section":    self._do_section,
                "bullets":    self._do_bullets,
                "statement":  self._do_statement,
                "comparison": self._do_comparison,
                "big_number": self._do_big_number,
                "code":       self._do_code,
                "end":        self._do_end,
            }.get(layout, self._do_bullets)
            fn(slide, s)

            # page number (skip first & last)
            if layout not in ("title", "end"):
                self._page_num(slide, i + 1, n)

        self.prs.save(out)
        return out

    # ── Helpers ─────────────────────────────────────────────

    def _c(self, hex_key: str) -> RGBColor:
        """Resolve a palette color key to RGBColor."""
        v = self.pal.get(hex_key, hex_key)
        if v.startswith("#"):
            v = v.lstrip("#")
            return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:], 16))
        return RGBColor(0, 0, 0)

    def _bg(self, slide, hex_key: str):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self._c(hex_key)

    def _rect(self, slide, pos, size, color_key: str, radius=False):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = slide.shapes.add_shape(
            shape_type,
            Inches(pos[0] * SW), Inches(pos[1] * SH),
            Inches(size[0] * SW), Inches(size[1] * SH),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._c(color_key)
        shape.line.fill.background()
        return shape

    def _oval(self, slide, pos, size, color_key: str):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(pos[0] * SW), Inches(pos[1] * SH),
            Inches(size[0] * SW), Inches(size[1] * SH),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._c(color_key)
        shape.line.fill.background()
        return shape

    def _tb(self, slide, text, spec_key, *,
            sz=18, bold=False, italic=False, color="text",
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, wrap=True,
            line_sp=None, font_name="Microsoft YaHei"):
        """Add a positioned textbox."""
        spec = LAYOUTS.get(spec_key, {})  # not used directly; caller passes lkey
        # This is a lower-level helper; see _txt below.
        pass

    def _txt(self, slide, text, layout_key, spec_key, *,
             sz=18, bold=False, italic=False, color="text",
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, wrap=True,
             line_sp=None, font_name="Microsoft YaHei"):
        """Add text at a layout-defined position."""
        lk = LAYOUTS.get(layout_key, {})
        sp = lk.get(spec_key, {"p": (0.1, 0.1), "s": (0.8, 0.1)})
        tb = slide.shapes.add_textbox(
            Inches(sp["p"][0] * SW), Inches(sp["p"][1] * SH),
            Inches(sp["s"][0] * SW), Inches(sp["s"][1] * SH),
        )
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.auto_size = None
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        if line_sp:
            p.line_spacing = Pt(line_sp)
        run = p.runs[0] if p.runs else p.add_run()
        if not p.runs[0].text and text:
            p.runs[0].text = text
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = self._c(color)
        run.font.name = font_name
        return tb

    def _para(self, tf, text, *, sz=18, bold=False, italic=False,
              color="text", align=PP_ALIGN.LEFT, line_sp=None,
              space_after=None, font_name="Microsoft YaHei"):
        """Append a paragraph to an existing text frame."""
        p = tf.add_paragraph()
        p.text = text
        p.alignment = align
        if line_sp:
            p.line_spacing = Pt(line_sp)
        if space_after is not None:
            p.space_after = Pt(space_after)
        r = p.runs[0] if p.runs else p.add_run()
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = self._c(color)
        r.font.name = font_name
        return p

    def _page_num(self, slide, num, total):
        tb = slide.shapes.add_textbox(
            Inches(0.88 * SW), Inches(0.92 * SH),
            Inches(0.08 * SW), Inches(0.04 * SH),
        )
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = f"{num} / {total}"
        r.font.size = Pt(10)
        r.font.color.rgb = self._c("text3")
        r.font.name = "Microsoft YaHei"

    # ── Layout Renderers ────────────────────────────────────

    def _do_title(self, slide, s):
        # Decorative circles
        self._oval(slide, (0.82, 0.68), (0.16, 0.28), "accent1")
        self._oval(slide, (0.89, 0.55), (0.10, 0.18), "accent2")
        # Left accent bar
        self._rect(slide, (0.04, 0.28), (0.008, 0.22), "accent1")
        # Title
        self._txt(slide, s.get("title", ""), "title", "title",
                  sz=40, bold=True, align=PP_ALIGN.LEFT, color="text")
        # Subtitle
        sub = s.get("subtitle", s.get("body_text", ""))
        self._txt(slide, sub, "title", "subtitle",
                  sz=22, align=PP_ALIGN.LEFT, color="text2")

    def _do_section(self, slide, s):
        # Full accent background
        self._bg(slide, "accent1")
        # Decorative circle (subtle)
        self._oval(slide, (0.72, 0.08), (0.32, 0.56), "accent2")
        # Large section number
        num = str(s.get("act", s.get("slide_num", "")))
        if len(num) == 1:
            num = f"0{num}"
        lk = LAYOUTS["section"]
        sp = lk["number"]
        tb = slide.shapes.add_textbox(
            Inches(sp["p"][0] * SW), Inches(sp["p"][1] * SH),
            Inches(sp["s"][0] * SW), Inches(sp["s"][1] * SH),
        )
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = num
        r.font.size = Pt(72)
        r.font.bold = True
        r.font.color.rgb = self._c("on_accent")
        r.font.name = "Impact"
        # Title (white)
        self._txt(slide, s.get("title", ""), "section", "title",
                  sz=36, bold=True, color="on_accent")
        # Subtitle
        sub = s.get("subtitle", s.get("body_text", ""))
        if sub:
            self._txt(slide, sub, "section", "subtitle",
                      sz=18, color="on_accent")

    def _do_bullets(self, slide, s):
        # Left accent bar
        self._rect(slide, (0.05, 0.06), (0.006, 0.11), "accent1")
        # Title
        self._txt(slide, s.get("title", ""), "bullets", "title",
                  sz=30, bold=True, color="text")
        # Bullets
        bullets = s.get("bullets", [])
        body = s.get("body_text", "")
        lk = LAYOUTS["bullets"]
        sp = lk["body"]
        tb = slide.shapes.add_textbox(
            Inches(sp["p"][0] * SW), Inches(sp["p"][1] * SH),
            Inches(sp["s"][0] * SW), Inches(sp["s"][1] * SH),
        )
        tf = tb.text_frame
        tf.word_wrap = True

        first = True
        for b in bullets:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(10)
            # Accent bullet character
            r1 = p.add_run()
            r1.text = "●  "
            r1.font.size = Pt(10)
            r1.font.color.rgb = self._c("accent1")
            r1.font.name = "Microsoft YaHei"
            # Bullet text
            r2 = p.add_run()
            r2.text = b
            r2.font.size = Pt(18)
            r2.font.color.rgb = self._c("text")
            r2.font.name = "Microsoft YaHei"

        # Body text below bullets
        if body:
            if first:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(14)
            r = p.add_run()
            r.text = body
            r.font.size = Pt(16)
            r.font.color.rgb = self._c("text2")
            r.font.name = "Microsoft YaHei"

    def _do_statement(self, slide, s):
        # Surface background panel
        self._rect(slide, (0.06, 0.16), (0.88, 0.62), "surface", radius=True)
        # Left accent bar
        self._rect(slide, (0.07, 0.20), (0.006, 0.52), "accent1")
        # Quote mark decoration
        lk = LAYOUTS["statement"]
        sp_q = lk["quote"]
        tb = slide.shapes.add_textbox(
            Inches((sp_q["p"][0] + 0.04) * SW), Inches(sp_q["p"][1] * SH),
            Inches(sp_q["s"][0] * SW), Inches(sp_q["s"][1] * SH),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        # Opening quote mark
        r0 = p.add_run()
        r0.text = "\u201c"
        r0.font.size = Pt(56)
        r0.font.color.rgb = self._c("accent1")
        r0.font.name = "Georgia"
        # Main text
        r1 = p.add_run()
        r1.text = f"\n{s.get('title', s.get('body_text', ''))}"
        r1.font.size = Pt(28)
        r1.font.bold = True
        r1.font.color.rgb = self._c("text")
        r1.font.name = "Microsoft YaHei"
        p.line_spacing = Pt(42)
        # Attribution
        attr = s.get("subtitle", "")
        if attr:
            self._txt(slide, f"— {attr}", "statement", "attr",
                      sz=16, italic=True, color="text3", align=PP_ALIGN.LEFT)

    def _do_comparison(self, slide, s):
        # Title
        self._txt(slide, s.get("title", ""), "comparison", "title",
                  sz=30, bold=True, color="text")
        # Center divider
        self._rect(slide, (0.497, 0.20), (0.004, 0.66), "divider")
        # Left column header bar
        self._rect(slide, (0.08, 0.20), (0.40, 0.006), "accent1")
        # Right column header bar
        self._rect(slide, (0.52, 0.20), (0.40, 0.006), "accent2")

        lk = LAYOUTS["comparison"]
        # Left column
        left_items = s.get("left_bullets", s.get("bullets", []))
        left_title = s.get("left_title", "")
        sp_l = lk["left"]
        tb_l = slide.shapes.add_textbox(
            Inches(sp_l["p"][0] * SW), Inches(sp_l["p"][1] * SH),
            Inches(sp_l["s"][0] * SW), Inches(sp_l["s"][1] * SH),
        )
        tf_l = tb_l.text_frame
        tf_l.word_wrap = True
        if left_title:
            p = tf_l.paragraphs[0]
            r = p.add_run()
            r.text = left_title
            r.font.size = Pt(20)
            r.font.bold = True
            r.font.color.rgb = self._c("accent1")
            r.font.name = "Microsoft YaHei"
            p.space_after = Pt(12)
        for i, item in enumerate(left_items):
            p = tf_l.add_paragraph() if i > 0 or left_title else tf_l.paragraphs[0]
            r = p.add_run()
            r.text = f"●  {item}"
            r.font.size = Pt(16)
            r.font.color.rgb = self._c("text")
            r.font.name = "Microsoft YaHei"
            p.space_after = Pt(8)

        # Right column
        right_items = s.get("right_bullets", [])
        right_title = s.get("right_title", "")
        sp_r = lk["right"]
        tb_r = slide.shapes.add_textbox(
            Inches(sp_r["p"][0] * SW), Inches(sp_r["p"][1] * SH),
            Inches(sp_r["s"][0] * SW), Inches(sp_r["s"][1] * SH),
        )
        tf_r = tb_r.text_frame
        tf_r.word_wrap = True
        if right_title:
            p = tf_r.paragraphs[0]
            r = p.add_run()
            r.text = right_title
            r.font.size = Pt(20)
            r.font.bold = True
            r.font.color.rgb = self._c("accent2")
            r.font.name = "Microsoft YaHei"
            p.space_after = Pt(12)
        for i, item in enumerate(right_items):
            p = tf_r.add_paragraph() if i > 0 or right_title else tf_r.paragraphs[0]
            r = p.add_run()
            r.text = f"●  {item}"
            r.font.size = Pt(16)
            r.font.color.rgb = self._c("text")
            r.font.name = "Microsoft YaHei"
            p.space_after = Pt(8)

    def _do_big_number(self, slide, s):
        # Big number
        lk = LAYOUTS["big_number"]
        sp = lk["number"]
        tb = slide.shapes.add_textbox(
            Inches(sp["p"][0] * SW), Inches(sp["p"][1] * SH),
            Inches(sp["s"][0] * SW), Inches(sp["s"][1] * SH),
        )
        p = tb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = s.get("highlight", s.get("title", ""))
        r.font.size = Pt(72)
        r.font.bold = True
        r.font.color.rgb = self._c("accent1")
        r.font.name = "Impact"
        # Label
        self._txt(slide, s.get("title", ""), "big_number", "label",
                  sz=22, bold=True, color="text", align=PP_ALIGN.CENTER)
        # Context
        ctx = s.get("body_text", s.get("subtitle", ""))
        if ctx:
            self._txt(slide, ctx, "big_number", "context",
                      sz=16, color="text2", align=PP_ALIGN.CENTER)

    def _do_code(self, slide, s):
        # Title
        self._txt(slide, s.get("title", ""), "code", "title",
                  sz=30, bold=True, color="text")
        # Code block background
        lk = LAYOUTS["code"]
        sp_c = lk["code"]
        code_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(sp_c["p"][0] * SW), Inches(sp_c["p"][1] * SH),
            Inches(sp_c["s"][0] * SW), Inches(sp_c["s"][1] * SH),
        )
        code_bg.fill.solid()
        code_bg.fill.fore_color.rgb = self._c("code_bg")
        code_bg.line.fill.background()
        # Code text
        code_text = s.get("code", s.get("body_text", ""))
        tb = slide.shapes.add_textbox(
            Inches((sp_c["p"][0] + 0.02) * SW), Inches((sp_c["p"][1] + 0.03) * SH),
            Inches((sp_c["s"][0] - 0.04) * SW), Inches((sp_c["s"][1] - 0.06) * SH),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        for i, line in enumerate(code_text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = line
            r.font.size = Pt(13)
            r.font.color.rgb = RGBColor(0xD4, 0xD4, 0xD4)
            r.font.name = "Consolas"
            p.line_spacing = Pt(20)
        # Annotations on the right
        notes = s.get("annotations", s.get("bullets", []))
        if notes:
            sp_n = lk["notes"]
            tb_n = slide.shapes.add_textbox(
                Inches(sp_n["p"][0] * SW), Inches(sp_n["p"][1] * SH),
                Inches(sp_n["s"][0] * SW), Inches(sp_n["s"][1] * SH),
            )
            tf_n = tb_n.text_frame
            tf_n.word_wrap = True
            for i, note in enumerate(notes):
                p = tf_n.paragraphs[0] if i == 0 else tf_n.add_paragraph()
                r = p.add_run()
                r.text = f"→  {note}"
                r.font.size = Pt(14)
                r.font.color.rgb = self._c("text2")
                r.font.name = "Microsoft YaHei"
                p.space_after = Pt(10)

    def _do_end(self, slide, s):
        # Full accent background
        self._bg(slide, "accent1")
        # Decorative circles
        self._oval(slide, (0.02, 0.04), (0.14, 0.25), "accent2")
        self._oval(slide, (0.88, 0.72), (0.12, 0.22), "accent2")
        # Thanks
        thanks = s.get("title", "谢谢")
        self._txt(slide, thanks, "end", "thanks",
                  sz=40, bold=True, color="on_accent", align=PP_ALIGN.CENTER)
        # Contact
        contact = s.get("body_text", s.get("subtitle", ""))
        if contact:
            self._txt(slide, contact, "end", "contact",
                      sz=18, color="on_accent", align=PP_ALIGN.CENTER)
