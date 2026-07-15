#!/usr/bin/env python3
"""
Comic PPT Renderer — Generate cartoon comic PPT with AI illustrations.

Uses python-pptx to create a comic book style presentation with:
  - Full-bleed cartoon illustrations per page
  - Speech bubbles with character dialogue
  - Narration bars with story text
  - Sound effects as comic-style bold text
  - Cover and ending slides with special layouts
"""

import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# ── Palette (festive warm tones, adapts to comic palette) ──

PALETTES = {
    "warm_pastel": {
        "primary": RGBColor(0xE8, 0x6F, 0xA0),   # pink
        "secondary": RGBColor(0xFF, 0xB7, 0x4D),  # warm orange
        "bg_dark": RGBColor(0x3D, 0x2C, 0x5E),    # deep purple
        "accent": RGBColor(0xFF, 0x7E, 0xB3),
    },
    "bright_fun": {
        "primary": RGBColor(0x4E, 0xA8, 0xDB),    # sky blue
        "secondary": RGBColor(0xFF, 0xD7, 0x00),  # gold
        "bg_dark": RGBColor(0x1A, 0x3C, 0x5E),    # navy
        "accent": RGBColor(0x4E, 0xA8, 0xDB),
    },
    "soft_edu": {
        "primary": RGBColor(0x5B, 0xB5, 0x8A),    # sage green
        "secondary": RGBColor(0x7E, 0xB8, 0xDA),  # soft blue
        "bg_dark": RGBColor(0x2B, 0x3E, 0x50),    # dark teal
        "accent": RGBColor(0x5B, 0xB5, 0x8A),
    },
    "festive": {
        "primary": RGBColor(0xFF, 0x7E, 0xB3),    # pink
        "secondary": RGBColor(0xFF, 0xD7, 0x00),  # gold
        "bg_dark": RGBColor(0x2D, 0x1B, 0x4E),    # deep purple
        "accent": RGBColor(0xFF, 0x7E, 0xB3),
    },
}

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x3A)
MUTED = RGBColor(0xAA, 0x99, 0xBB)
SFX_COLORS = [
    RGBColor(0xFF, 0x8C, 0x00),  # orange
    RGBColor(0xFF, 0x69, 0xB4),  # hot pink
    RGBColor(0xFF, 0xD7, 0x00),  # gold
]

FONT = "PingFang SC"
FONT_EA = "PingFang SC"

# Fallback font if PingFang not available
FONT_FALLBACKS = ["Microsoft YaHei", "Source Han Sans SC", "Noto Sans CJK SC", "SimHei"]


def set_cjk_font(run, font_name=FONT, ea_name=FONT_EA, bold=False):
    """Set CJK-compatible fonts on a text run."""
    run.font.name = font_name
    run.font.bold = bold
    rPr = run._r.get_or_add_rPr()
    ea_el = rPr.find(qn("a:ea"))
    if ea_el is None:
        ea_el = rPr.makeelement(qn("a:ea"), {"typeface": ea_name})
        rPr.append(ea_el)
    else:
        ea_el.set("typeface", ea_name)


def set_shape_alpha(shape, alpha_pct):
    """Set fill transparency. alpha_pct: 0=transparent, 100=opaque."""
    spPr = shape._element.spPr
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is None:
        return
    srgbClr = solidFill.find(qn("a:srgbClr"))
    if srgbClr is not None:
        for existing in srgbClr.findall(qn("a:alpha")):
            srgbClr.remove(existing)
        alpha_el = srgbClr.makeelement(qn("a:alpha"), {"val": str(int(alpha_pct * 1000))})
        srgbClr.append(alpha_el)


class ComicPPTRenderer:
    """Render comic pages as a PPT with cartoon illustrations."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.pal = PALETTES["festive"]
        self.total_pages = 0

    def render(self, pages, palette_name, plan_data, image_paths, out):
        """Render comic to PPT file."""
        self.pal = PALETTES.get(palette_name, PALETTES["festive"])
        self.total_pages = len(pages)

        for i, page in enumerate(pages):
            img_path = image_paths[i] if i < len(image_paths) else ""
            ptype = page.get("page_type", "scene")

            if ptype == "cover":
                self._build_cover(page, img_path, plan_data)
            elif ptype == "ending":
                self._build_ending(page, img_path, plan_data)
            else:
                self._build_story_page(page, img_path, i)

        os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
        self._save(out)
        return out

    # ── Text helpers ──────────────────────────────────

    def _add_text(self, slide, text, x, y, w, h, *,
                  size=14, bold=False, color=None, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.04)
        tf.margin_bottom = Inches(0.04)
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        if color:
            r.font.color.rgb = color
        set_cjk_font(r, bold=bold)
        return tb

    def _add_bubble(self, slide, text, x, y, w, h, char_name="", font_size=13):
        """Add speech bubble (rounded rect + triangle tail)."""
        bubble = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = WHITE
        bubble.line.color.rgb = RGBColor(0xE0, 0xD0, 0xC0)
        bubble.line.width = Pt(1.5)
        bubble.adjustments[0] = 0.2

        tf = bubble.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.06)
        tf.margin_bottom = Inches(0.06)

        if char_name:
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = char_name
            r.font.size = Pt(9)
            r.font.color.rgb = self.pal["accent"]
            r.font.bold = True
            set_cjk_font(r)
            p_text = tf.add_paragraph()
        else:
            p_text = tf.paragraphs[0]

        p_text.alignment = PP_ALIGN.LEFT
        r = p_text.add_run()
        r.text = text
        r.font.size = Pt(font_size)
        r.font.color.rgb = DARK_TEXT
        set_cjk_font(r)

        # Tail triangle
        tail = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(x + w * 0.3), Inches(y + h - 0.01),
            Inches(0.15), Inches(0.15)
        )
        tail.fill.solid()
        tail.fill.fore_color.rgb = WHITE
        tail.line.fill.background()
        tail.rotation = 180.0

    def _add_sfx(self, slide, text, x, y, size=36, color_idx=0):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3.0), Inches(1.0))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.italic = True
        r.font.color.rgb = SFX_COLORS[color_idx % len(SFX_COLORS)]
        set_cjk_font(r, bold=True)

    def _add_narration(self, slide, text, y=5.85, h=1.65):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(y), Inches(13.333), Inches(h)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.pal["bg_dark"]
        bar.line.fill.background()
        set_shape_alpha(bar, 75)

        self._add_text(slide, text, 0.6, y + 0.15, 12.1, h - 0.3,
                       size=15, color=WHITE)

    def _add_page_num(self, slide, num):
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(12.4), Inches(0.15), Inches(0.75), Inches(0.4)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = self.pal["accent"]
        badge.line.fill.background()
        badge.adjustments[0] = 0.5
        tf = badge.text_frame
        tf.word_wrap = False
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"{num}/{self.total_pages}"
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = WHITE
        set_cjk_font(r)

    def _add_image(self, slide, img_path, x=0, y=0, w=13.333, h=7.5):
        """Add image, or colored placeholder if missing."""
        if img_path and os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(x), Inches(y),
                                     width=Inches(w), height=Inches(h))
        else:
            # Colored placeholder
            rect = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(h)
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = self.pal["bg_dark"]
            rect.line.fill.background()

    # ── Page builders ─────────────────────────────────

    def _build_cover(self, page, img_path, plan_data):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_image(s, img_path)

        # Dark overlay at bottom
        overlay = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(4.5), Inches(13.333), Inches(3.0)
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = self.pal["bg_dark"]
        overlay.line.fill.background()
        set_shape_alpha(overlay, 70)

        title = plan_data.get("title", page.get("title", ""))
        self._add_text(s, title, 1.0, 5.0, 11.3, 1.2,
                       size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        moral = plan_data.get("moral", "")
        if moral:
            self._add_text(s, moral, 1.0, 6.2, 11.3, 0.6,
                           size=22, color=RGBColor(0xFF, 0xD7, 0x00),
                           align=PP_ALIGN.CENTER)

    def _build_story_page(self, page, img_path, slide_idx):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Full-bleed image (top 78%)
        self._add_image(s, img_path, y=0, h=5.8)

        # Page title badge
        title = page.get("title", "")
        if title:
            badge = s.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(0.3), Inches(0.2), Inches(2.8), Inches(0.55)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = self.pal["accent"]
            badge.line.fill.background()
            badge.adjustments[0] = 0.4
            tf = badge.text_frame
            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = title
            r.font.size = Pt(16)
            r.font.bold = True
            r.font.color.rgb = WHITE
            set_cjk_font(r)

        # Page number
        self._add_page_num(s, page.get("page_num", slide_idx + 1))

        # Speech bubbles
        dialogues = page.get("dialogue", [])
        if dialogues:
            bx, bw = 7.5, 4.8
            bh = 0.85 if len(dialogues) <= 2 else 0.72
            for i, dlg in enumerate(dialogues[:3]):
                by = 0.25 + i * (bh + 0.12)
                # Clean emoji from character name
                name = dlg.get("character", "")
                clean = "".join(c for c in name if ord(c) < 0x1F000 or ord(c) > 0x1FFFF).strip()
                self._add_bubble(s, dlg.get("text", ""),
                                 bx, by, bw, bh,
                                 char_name=clean, font_size=14)

        # Sound effect
        sfx = page.get("sound_effect", "")
        if sfx:
            self._add_sfx(s, sfx, 4.5, 3.8, size=42, color_idx=slide_idx)

        # Narration
        story = page.get("story_text", "")
        if story:
            self._add_narration(s, story)

        # Mood
        mood = page.get("mood", "")
        if mood:
            self._add_text(s, f"· {mood} ·", 10.5, 7.0, 2.5, 0.4,
                           size=10, color=MUTED, align=PP_ALIGN.RIGHT)

    def _build_ending(self, page, img_path, plan_data):
        s = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_image(s, img_path)

        # Dark overlay (semi-transparent)
        overlay = s.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(13.333), Inches(7.5)
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = self.pal["bg_dark"]
        overlay.line.fill.background()
        set_shape_alpha(overlay, 45)

        # Story text
        story = page.get("story_text", "")
        if story:
            self._add_text(s, story, 1.5, 1.5, 10.3, 1.5,
                           size=20, color=WHITE, align=PP_ALIGN.CENTER)

        # Dialogue
        dialogues = page.get("dialogue", [])
        if dialogues:
            y = 3.2
            for dlg in dialogues[:2]:
                name = dlg.get("character", "")
                clean = "".join(c for c in name if ord(c) < 0x1F000 or ord(c) > 0x1FFFF).strip()
                line = f"{clean}：{dlg.get('text', '')}" if clean else dlg.get("text", "")
                self._add_text(s, line, 2.0, y, 9.3, 0.6,
                               size=16, color=WHITE, align=PP_ALIGN.CENTER)
                y += 0.55

        # Moral
        moral = plan_data.get("moral", "故事结束")
        self._add_text(s, moral, 1.0, 5.2, 11.3, 1.0,
                       size=36, bold=True,
                       color=RGBColor(0xFF, 0xD7, 0x00),
                       align=PP_ALIGN.CENTER)

        # Page number
        self._add_page_num(s, self.total_pages)

    # ── Save ──────────────────────────────────────────

    def _save(self, path):
        """Save PPTX and strip template thumbnail."""
        self.prs.save(path)
        import zipfile, re
        tmp = path + ".tmp"
        try:
            with zipfile.ZipFile(path, "r") as zin, zipfile.ZipFile(tmp, "w") as zout:
                for item in zin.infolist():
                    if item.filename.lower().startswith("docprops/thumbnail"):
                        continue
                    data = zin.read(item.filename)
                    if item.filename == "[Content_Types].xml":
                        content = data.decode("utf-8")
                        content = re.sub(
                            r'<Override[^>]*PartName="/docProps/thumbnail[^"]*"[^>]*/?>',
                            '', content, flags=re.IGNORECASE)
                        data = content.encode("utf-8")
                    if item.filename == "_rels/.rels":
                        content = data.decode("utf-8")
                        content = re.sub(
                            r'<Relationship[^>]*Target="docProps/thumbnail[^"]*"[^>]*/?>',
                            '', content, flags=re.IGNORECASE)
                        data = content.encode("utf-8")
                    zout.writestr(item, data)
            os.replace(tmp, path)
        except Exception as e:
            print(f"  Warning: thumbnail cleanup: {e}")
            if os.path.exists(tmp):
                os.remove(tmp)
